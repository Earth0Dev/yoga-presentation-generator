from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import random
import json
import requests
from datetime import datetime
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
import io

app = Flask(
    __name__,
    static_folder=".",      # serve files from the same folder
    template_folder="."     # allow index.html to be served
)

CORS(app)

@app.route("/")
def index():
    return app.send_static_file("index.html")


# API Keys
UNSPLASH_API_KEY = os.getenv('UNSPLASH_API_KEY', 'UNSPLASH_API_KEY')
OPENROUTER_API_KEY = os.getenv('OPENROUTER_API_KEY', 'OPENROUTER_API_KEY')

# Clean yoga names
YOGA_ASANAS = [
    "Tadasana", "Ardha Chakrasana", "Kati Chakrasana", "Dandasana", "Bhadrasana",
    "Padmasana", "Vajrasana", "Utthan Mandukasana", "Kakasana", "Parvatasana",
    "Makarasana", "Uttanpadasana", "Setu Bandhasana", "Viparita Karani", 
    "Saral Matsyasana", "Shavasana", "Pranayama", "Yoga Mudras", "Meditation"
]

# Enhanced Canva-style color palettes
CANVA_PALETTES = [
    {  # Modern Wellness
        "primary": RGBColor(74, 144, 226),    # Calm Blue
        "secondary": RGBColor(76, 175, 80),   # Fresh Green
        "accent": RGBColor(255, 193, 7),      # Warm Yellow
        "background": RGBColor(250, 250, 250),# Pure White
        "text": RGBColor(33, 33, 33),         # Dark Charcoal
        "highlight": RGBColor(156, 39, 176),  # Spiritual Purple
        "success": RGBColor(56, 142, 60),     # Deep Green
        "warning": RGBColor(255, 152, 0),     # Orange
        "light_bg": RGBColor(245, 248, 250)   # Light Blue Gray
    },
    {  # Earthy Yoga
        "primary": RGBColor(121, 85, 72),     # Earth Brown
        "secondary": RGBColor(104, 159, 56),  # Nature Green
        "accent": RGBColor(255, 167, 38),     # Sunset Orange
        "background": RGBColor(253, 249, 236),# Cream
        "text": RGBColor(55, 45, 45),         # Dark Brown
        "highlight": RGBColor(175, 180, 43),  # Olive
        "success": RGBColor(67, 160, 71),     # Forest Green
        "warning": RGBColor(229, 115, 115),   # Clay Red
        "light_bg": RGBColor(252, 243, 231)   # Light Beige
    }
]

def get_comprehensive_yoga_content(asanas_name, student_name=""):
    """Get comprehensive AI content with proper structure"""
    try:
        print(f"🔍 Fetching structured AI content for: {asanas_name}")
        
        prompt = f"""
        Create a comprehensive yoga presentation for {asanas_name} with exactly 13 slides.
        Each slide should have clear structure with headings, subheadings, and 5-6 detailed bullet points.
        
        FORMAT EACH SLIDE LIKE THIS:
        
        SLIDE 1: MASTERING {asanas_name.upper()}
        🎯 SUBHEADING: Transform Your Yoga Practice
        • Foundational standing pose suitable for all experience levels
        • Builds core strength, improves balance and body awareness
        • Perfect for morning routines and daily practice sessions
        • Creates deep connection between body, mind and breathing
        • Enhances overall posture and spinal alignment
        • Recommended by yoga therapists for holistic wellness
        
        SLIDE 2: STEP-BY-STEP GUIDANCE
        📝 SUBHEADING: Perfect Your Alignment & Form
        • Start with feet hip-width apart, weight evenly distributed
        • Ground through all four corners of each foot firmly
        • Engage thigh muscles while gently lifting the kneecaps
        • Lengthen spine upward, creating space between vertebrae
        • Relax shoulders away from ears, broaden collarbones
        • Maintain steady gaze forward with soft focus
        
        SLIDE 3: BREATHING TECHNIQUES
        🌬️ SUBHEADING: Sync Movement with Conscious Breath
        • Practice deep diaphragmatic breathing throughout the pose
        • Inhale deeply to lengthen spine and expand chest fully
        • Exhale completely to ground feet and stabilize position
        • Maintain steady, rhythmic breathing pattern consistently
        • Coordinate breath with subtle micro-movements in pose
        • Use ujjayi breathing for enhanced focus and warmth
        
        Continue with 10 more slides covering benefits, modifications, safety, etc.
        Each slide must have: MAIN TITLE, EMOJI SUBHEADING, and 5-6 DETAILED BULLET POINTS.
        Make content comprehensive and practical for yoga practitioners.
        """
        
        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json"
        }
        
        payload = {
            "model": "openai/gpt-3.5-turbo",
            "messages": [
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            "max_tokens": 3000,
            "temperature": 0.7
        }
        
        response = requests.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers=headers,
            json=payload,
            timeout=30
        )
        
        if response.status_code == 200:
            result = response.json()
            content = result['choices'][0]['message']['content']
            return parse_structured_content(content, asanas_name)
        else:
            raise Exception(f"OpenRouter API error: {response.status_code}")
            
    except Exception as e:
        print(f"🚨 AI Content Error: {e}")
        return get_enhanced_fallback_content(asanas_name)

def parse_structured_content(content, asanas_name):
    """Parse AI content into structured slides with titles and subheadings"""
    slides = []
    lines = content.split('\n')
    
    current_slide = {"title": "", "subheading": "", "content": []}
    
    for line in lines:
        line = line.strip()
        
        if line.startswith('SLIDE') or (line.isupper() and len(line) > 10):
            # New slide title
            if current_slide["title"]:
                slides.append(format_slide_content(current_slide))
            current_slide = {"title": line, "subheading": "", "content": []}
            
        elif 'SUBHEADING:' in line or any(emoji in line for emoji in ['🎯', '📝', '🌬️', '💪', '🧠', '🌟', '🚀', '⚠️', '📅', '👨‍🏫', '📚', '💼', '🎉']):
            # Subheading line
            current_slide["subheading"] = line
            
        elif line.startswith('•') or line.startswith('-'):
            # Content bullet point
            current_slide["content"].append(line)
    
    # Add the last slide
    if current_slide["title"]:
        slides.append(format_slide_content(current_slide))
    
    # Ensure we have exactly 13 slides with enhanced content
    enhanced_slides = get_enhanced_fallback_content(asanas_name)
    while len(slides) < 13:
        slides.append(enhanced_slides[len(slides) % len(enhanced_slides)])
    
    return slides[:13]

def format_slide_content(slide_data):
    """Format slide content with proper structure"""
    title = slide_data.get("title", "Yoga Practice")
    subheading = slide_data.get("subheading", "")
    content_lines = slide_data.get("content", [])
    
    # Ensure we have 5-6 bullet points for comprehensive content
    while len(content_lines) < 5:
        content_lines.append("• Additional detailed point for comprehensive coverage")
    
    # Limit to 6 bullet points maximum
    content = '\n'.join(content_lines[:6])
    
    # Combine all elements
    full_content = f"{title}\n{subheading}\n{content}"
    return truncate_structured_content(full_content)

def truncate_structured_content(content, max_chars=500):
    """Ensure structured content fits without overflow"""
    if len(content) > max_chars:
        lines = content.split('\n')
        truncated = []
        total_length = 0
        
        for line in lines:
            if total_length + len(line) < max_chars:
                truncated.append(line)
                total_length += len(line) + 1
            else:
                break
        
        result = '\n'.join(truncated)
        return result if result.strip() else content[:max_chars-3] + '...'
    return content

def get_enhanced_fallback_content(asanas_name):
    """Enhanced fallback content with more bullet points"""
    return [
        f"MASTERING {asanas_name.upper()}\n🎯 SUBHEADING: Transform Your Yoga Journey\n• Foundational pose suitable for beginners to advanced practitioners\n• Builds comprehensive strength, balance and body awareness\n• Perfect integration into daily morning or evening routines\n• Creates profound connection between physical and mental states\n• Significantly enhances overall posture and spinal health\n• Recommended by yoga therapists for holistic wellness approach",
        
        f"STEP-BY-STEP GUIDANCE\n📝 SUBHEADING: Perfect Your Alignment & Technique\n• Start with proper foot placement and weight distribution\n• Engage core muscles while maintaining relaxed breathing\n• Align spine vertically with natural curvature maintained\n• Position shoulders correctly away from ear position\n• Coordinate subtle movements with breath patterns\n• Maintain steady gaze and focused mental attention",
        
        f"BREATHING TECHNIQUES\n🌬️ SUBHEADING: Master Conscious Breathing Patterns\n• Practice deep diaphragmatic breathing throughout entire pose\n• Coordinate inhalation with expansion and lifting movements\n• Synchronize exhalation with grounding and stabilizing actions\n• Maintain consistent, rhythmic breathing pattern always\n• Incorporate advanced pranayama techniques when ready\n• Use breath as anchor for mental focus and concentration",
        
        f"PHYSICAL BENEFITS\n💪 SUBHEADING: Transform Your Body Completely\n• Dramatically improves flexibility and joint mobility\n• Significantly strengthens core and postural muscles\n• Enhances overall body alignment and spinal health\n• Increases blood circulation and oxygen flow\n• Improves balance and proprioception significantly\n• Reduces risk of injury and chronic pain",
        
        f"MENTAL BENEFITS\n🧠 SUBHEADING: Achieve Mental Clarity & Peace\n• Effectively reduces stress and anxiety levels\n• Improves mental focus and concentration abilities\n• Enhances mind-body connection and self-awareness\n• Promotes emotional balance and stability\n• Increases mindfulness and present moment awareness\n• Supports overall mental health and wellbeing",
        
        f"BEGINNER FRIENDLY\n🌟 SUBHEADING: Start Your Journey Confidently\n• Simple modifications available for all ability levels\n• Progressive learning path with clear milestones\n• Patient, supportive approach to skill development\n• Encouraging community and resources available\n• Safe practice guidelines for new practitioners\n• Celebrating small victories and progress",
        
        f"ADVANCED VARIATIONS\n🚀 SUBHEADING: Challenge Your Practice Further\n• Extended duration holds for strength building\n• Complex variations for experienced practitioners\n• Integration into flowing vinyasa sequences\n• Advanced breathing and bandha applications\n• Partner and assisted variations available\n• Creative expressions and personal adaptations",
        
        f"SAFETY PRECAUTIONS\n⚠️ SUBHEADING: Practice Smart & Stay Safe\n• Always listen to your body's signals and limitations\n• Avoid pushing beyond comfortable range of motion\n• Proper warm-up and preparation are absolutely essential\n• Consult healthcare providers for existing conditions\n• Use props and modifications when necessary\n• Practice under qualified guidance when starting",
        
        f"DAILY PRACTICE ROUTINE\n📅 SUBHEADING: Build Consistent Habits\n• Morning practice ideal for energy and focus\n• 15-20 minutes daily for optimal results\n• Gradual progression in difficulty and duration\n• Regular self-assessment and adjustment\n• Integration with other wellness practices\n• Tracking progress and celebrating improvements",
        
        f"TEACHING METHODOLOGY\n👨‍🏫 SUBHEADING: Share Knowledge Effectively\n• Clear, concise verbal cues and instructions\n• Comprehensive visual demonstrations and examples\n• Individualized adjustments and modifications\n• Positive, encouraging feedback and reinforcement\n• Safe and supportive learning environment\n• Progressive skill building approach",
        
        f"PHILOSOPHICAL FOUNDATIONS\n📚 SUBHEADING: Deepen Your Understanding\n• Ancient wisdom and modern science integration\n• Holistic approach to health and wellness\n• Connection to larger yoga philosophy system\n• Spiritual dimensions of physical practice\n• Ethical principles and lifestyle applications\n• Personal transformation through consistent practice",
        
        f"MODERN APPLICATIONS\n💼 SUBHEADING: Integrate Into Daily Life\n• Office chair variations for workplace wellness\n• Quick 5-minute break routines for busy schedules\n• Effective stress management tool for modern life\n• Family and group practice opportunities\n• Community building through shared practice\n• Lifestyle integration for sustainable benefits",
        
        f"CONCLUSION & NEXT STEPS\n🎉 SUBHEADING: Continue Your Growth Journey\n• Consistent daily practice is absolutely essential\n• Progressive learning path with clear milestones\n• Enjoy the process and celebrate each achievement\n• Share knowledge and experience with others\n• Explore related poses and deeper practices\n• Lifetime journey of learning and growth"
    ]

def get_canva_style_image(asanas_name, slide_number):
    """Get aesthetic yoga images matching slide content"""
    try:
        # Different aesthetic search terms for each slide type
        aesthetic_terms = [
            "minimalist yoga aesthetic modern", "yoga alignment professional detailed", 
            "breathing meditation peaceful serene", "yoga flexibility artistic beautiful",
            "mental wellness meditation calm", "beginner yoga gentle supportive",
            "advanced yoga challenging achievement", "yoga safety careful mindful",
            "daily routine morning yoga fresh", "yoga teaching instructor guidance",
            "spiritual yoga philosophy deep", "modern yoga lifestyle balanced",
            "yoga celebration achievement success"
        ]
        
        search_query = aesthetic_terms[slide_number % len(aesthetic_terms)]
        
        if UNSPLASH_API_KEY:
            url = "https://api.unsplash.com/search/photos"
            params = {
                'query': f"{search_query} wellness minimalist professional",
                'client_id': UNSPLASH_API_KEY,
                'per_page': 15,
                'orientation': 'landscape'
            }
            
            response = requests.get(url, params=params, timeout=15)
            
            if response.status_code == 200:
                data = response.json()
                if data['results']:
                    image_index = slide_number % len(data['results'])
                    image = data['results'][image_index]
                    return {
                        "url": image['urls']['regular'],
                        "alt": image['alt_description'] or f"Professional {search_query}",
                        "photographer": image['user']['name'],
                        "source": "unsplash"
                    }
        
        # Enhanced curated aesthetic fallback images
        aesthetic_images = [
            "https://images.unsplash.com/photo-1544367567-0f2fcb009e0b?ixlib=rb-4.0.3&w=800&fit=crop&q=80",
            "https://images.unsplash.com/photo-1506126613408-eca07ce68773?ixlib=rb-4.0.3&w=800&fit=crop&q=80", 
            "https://images.unsplash.com/photo-1575052814086-f385e2e2ad1b?ixlib=rb-4.0.3&w=800&fit=crop&q=80",
            "https://images.unsplash.com/photo-1518611012118-696072aa579a?ixlib=rb-4.0.3&w=800&fit=crop&q=80",
            "https://images.unsplash.com/photo-1545389336-cf090694435e?ixlib=rb-4.0.3&w=800&fit=crop&q=80",
            "https://images.unsplash.com/photo-1599901854545-de86350faf80?ixlib=rb-4.0.3&w=800&fit=crop&q=80",
            "https://images.unsplash.com/photo-1549576490-b0b4831ef60a?ixlib=rb-4.0.3&w=800&fit=crop&q=80",
            "https://images.unsplash.com/photo-1599447292183-5dc08a0a9763?ixlib=rb-4.0.3&w=800&fit=crop&q=80",
            "https://images.unsplash.com/photo-1500462918059-b1a0cb512f1d?ixlib=rb-4.0.3&w=800&fit=crop&q=80",
            "https://images.unsplash.com/photo-1575052814086-f385e2e2ad1b?ixlib=rb-4.0.3&w=800&fit=crop&q=80",
            "https://images.unsplash.com/photo-1518611012118-696072aa579a?ixlib=rb-4.0.3&w=800&fit=crop&q=80",
            "https://images.unsplash.com/photo-1545389336-cf090694435e?ixlib=rb-4.0.3&w=800&fit=crop&q=80",
            "https://images.unsplash.com/photo-1599901854545-de86350faf80?ixlib=rb-4.0.3&w=800&fit=crop&q=80"
        ]
        
        return {
            "url": aesthetic_images[slide_number % len(aesthetic_images)],
            "alt": f"Professional yoga - {search_query}",
            "photographer": "Wellness Photography Collective",
            "source": "curated"
        }
        
    except Exception as e:
        print(f"Image error: {e}")
        return None

def add_geometric_elements(slide, colors, slide_number):
    """Add Canva-style geometric design elements"""
    # Background pattern
    pattern = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    pattern.fill.solid()
    pattern.fill.fore_color.rgb = colors['light_bg']
    pattern.line.fill.background()
    
    # Enhanced decorative elements based on slide number
    circle_positions = [
        (0.3, 0.3), (12.7, 0.3), (0.3, 6.7), (12.7, 6.7),
        (2.5, 1.5), (10.5, 1.5), (2.5, 5.5), (10.5, 5.5),
        (6.5, 0.8), (6.5, 6.2)
    ]
    
    for i, (left, top) in enumerate(circle_positions[:6 + slide_number % 5]):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.2 + (i * 0.05)), Inches(0.2 + (i * 0.05))
        )
        circle.fill.solid()
        if i % 3 == 0:
            circle.fill.fore_color.rgb = colors['accent']
        elif i % 3 == 1:
            circle.fill.fore_color.rgb = colors['highlight'] 
        else:
            circle.fill.fore_color.rgb = colors['secondary']
        circle.line.fill.background()
        circle.fill.transparency = 0.6 + (i * 0.03)

def set_background_gradient(slide, colors):
    """Set beautiful gradient background"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[0].color.rgb = colors['light_bg']
    fill.gradient_stops[1].position = 1
    fill.gradient_stops[1].color.rgb = RGBColor(240, 245, 250)

def create_canva_style_powerpoint(presentation_data):
    """Create professional Canva-style PowerPoint with enhanced content"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    colors = random.choice(CANVA_PALETTES)
    
    # Enhanced cover slide with premium design
    cover_slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_gradient(cover_slide, colors)
    add_geometric_elements(cover_slide, colors, 0)
    
    # Enhanced cover title with multiple font sizes and styles
    title_box = cover_slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = presentation_data['asanas'][0].upper()
    
    # Main title - large and bold
    title_frame.paragraphs[0].font.size = Pt(58)
    title_frame.paragraphs[0].font.color.rgb = colors['primary']
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.name = 'Arial Black'
    
    # Enhanced subtitle
    subtitle_box = cover_slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(11.3), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_text = "COMPREHENSIVE YOGA PRACTICE GUIDE"
    if presentation_data.get('student_name'):
        subtitle_text += f" • PERSONALIZED FOR {presentation_data['student_name'].upper()}"
    subtitle_frame.text = subtitle_text
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = colors['secondary']
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.name = 'Arial'
    subtitle_frame.paragraphs[0].font.bold = True
    
    # Enhanced decorative element
    decor = cover_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.8), Inches(4.3), Inches(0.15))
    decor.fill.solid()
    decor.fill.fore_color.rgb = colors['accent']
    decor.line.fill.background()
    
    # Enhanced date and details
    date_box = cover_slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.3), Inches(0.6))
    date_frame = date_box.text_frame
    date_text = f"CREATED • {datetime.now().strftime('%B %d, %Y').upper()} • COMPREHENSIVE 13-SLIDE GUIDE"
    date_frame.text = date_text
    date_frame.paragraphs[0].font.size = Pt(16)
    date_frame.paragraphs[0].font.color.rgb = colors['text']
    date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    date_frame.paragraphs[0].font.name = 'Arial'
    date_frame.paragraphs[0].font.italic = True
    
    # Content slides with enhanced design and more content
    for i, slide_data in enumerate(presentation_data['slides']):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background_gradient(slide, colors)
        add_geometric_elements(slide, colors, i)
        
        # Parse structured content
        content_lines = slide_data['content'].split('\n')
        title = content_lines[0] if content_lines else f"Slide {i+1}"
        subheading = ""
        bullet_points = []
        
        for line in content_lines[1:]:
            if line.strip() and ('SUBHEADING:' in line or any(emoji in line for emoji in ['🎯', '📝', '🌬️', '💪', '🧠', '🌟', '🚀', '⚠️', '📅', '👨‍🏫', '📚', '💼', '🎉'])):
                subheading = line
            elif line.strip() and (line.startswith('•') or line.startswith('-')):
                bullet_points.append(line)
        
        # Ensure we have sufficient content (5-6 bullet points)
        while len(bullet_points) < 5:
            bullet_points.append("• Additional comprehensive guidance for complete practice")
        
        # Enhanced header with slide number
        header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9))
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = colors['primary']
        header_bg.line.fill.background()
        
        slide_no = slide.shapes.add_textbox(Inches(11.5), Inches(0.3), Inches(1.2), Inches(0.5))
        slide_no_frame = slide_no.text_frame
        slide_no_frame.text = f"SLIDE {i+1:02d}"
        slide_no_frame.paragraphs[0].font.size = Pt(18)
        slide_no_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        slide_no_frame.paragraphs[0].font.bold = True
        slide_no_frame.paragraphs[0].font.name = 'Arial Black'
        
        # Enhanced main title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(8), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title[:45] + ('...' if len(title) > 45 else '')
        title_frame.paragraphs[0].font.size = Pt(22)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.name = 'Arial'
        
        # Enhanced content area with modern card design
        content_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(6.5), Inches(5.2))
        content_card.fill.solid()
        content_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        content_card.line.color.rgb = colors['primary']
        content_card.line.width = Pt(2.5)
        content_card.shadow.inherit = False
        
        # Enhanced subheading (if exists)
        content_top = Inches(1.5)
        if subheading:
            subheading_box = slide.shapes.add_textbox(Inches(1), content_top, Inches(6), Inches(0.7))
            subheading_frame = subheading_box.text_frame
            subheading_frame.text = subheading
            subheading_frame.paragraphs[0].font.size = Pt(16)
            subheading_frame.paragraphs[0].font.color.rgb = colors['secondary']
            subheading_frame.paragraphs[0].font.bold = True
            subheading_frame.paragraphs[0].font.name = 'Arial'
            content_top += Inches(0.8)
        
        # Enhanced bullet points with more content
        bullet_content = '\n'.join(bullet_points[:6])  # Up to 6 bullet points
        content_box = slide.shapes.add_textbox(Inches(1), content_top, Inches(6), Inches(4.5))
        content_frame = content_box.text_frame
        content_frame.text = bullet_content
        content_frame.paragraphs[0].font.size = Pt(13)
        content_frame.paragraphs[0].font.color.rgb = colors['text']
        content_frame.paragraphs[0].font.name = 'Arial'
        
        # Enhanced styling for bullet points
        for j, paragraph in enumerate(content_frame.paragraphs):
            if j > 0:  # Skip first paragraph if it's empty
                paragraph.font.size = Pt(12)
                paragraph.level = 0
                if j % 2 == 0:
                    paragraph.font.color.rgb = colors['secondary']
                else:
                    paragraph.font.color.rgb = colors['text']
        
        # Enhanced image area with modern frame
        if slide_data.get('image'):
            try:
                img_stream = download_image(slide_data['image']['url'])
                if img_stream:
                    # Enhanced image container with shadow
                    img_container = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(1.3), Inches(4.4), Inches(4.2)
                    )
                    img_container.fill.solid()
                    img_container.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    img_container.line.color.rgb = colors['accent']
                    img_container.line.width = Pt(3.5)
                    
                    # Actual image
                    img = slide.shapes.add_picture(
                        img_stream, Inches(8), Inches(1.5), Inches(4), Inches(3.6)
                    )
                    
                    # Enhanced image caption
                    caption_box = slide.shapes.add_textbox(Inches(7.8), Inches(5.6), Inches(4.4), Inches(0.4))
                    caption_frame = caption_box.text_frame
                    caption_text = f"📸 PROFESSIONAL IMAGE • {slide_data['image']['photographer']}"
                    caption_frame.text = caption_text
                    caption_frame.paragraphs[0].font.size = Pt(10)
                    caption_frame.paragraphs[0].font.color.rgb = colors['text']
                    caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    caption_frame.paragraphs[0].font.name = 'Arial'
                    caption_frame.paragraphs[0].font.italic = True
            except Exception as e:
                print(f"Slide {i+1} image error: {e}")
                # Enhanced stylish placeholder with more content
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(1.3), Inches(4.4), Inches(4.2)
                )
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = colors['highlight']
                
                placeholder_text = slide.shapes.add_textbox(Inches(7.8), Inches(2.8), Inches(4.4), Inches(1.5))
                placeholder_frame = placeholder_text.text_frame
                placeholder_content = f"🎨\nVISUAL INSPIRATION\n{subheading.split(':')[-1] if ':' in subheading else 'Professional Yoga Practice'}\nComprehensive Guidance"
                placeholder_frame.text = placeholder_content
                placeholder_frame.paragraphs[0].font.size = Pt(16)
                placeholder_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                placeholder_frame.paragraphs[0].font.name = 'Arial'
                placeholder_frame.paragraphs[0].font.bold = True
        
        # Enhanced footer with decorative element
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.15))
        footer.fill.solid()
        footer.fill.fore_color.rgb = colors['accent']
        footer.line.fill.background()
    
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream

def download_image(url):
    """Download image from URL"""
    try:
        response = requests.get(url, timeout=15)
        if response.status_code == 200:
            return io.BytesIO(response.content)
    except:
        pass
    return None

@app.route('/generate', methods=['POST', 'OPTIONS'])
def generate_presentation():
    if request.method == 'OPTIONS':
        return '', 200
    
    try:
        data = request.get_json()
        asanas = data.get('asanas', [])
        student_name = data.get('student_name', '').strip()
        include_images = data.get('include_images', True)
        
        if not asanas:
            return jsonify({'status': 'error', 'message': 'Please select one asana'}), 400
        
        main_asana = asanas[0]
        
        # Get enhanced structured AI content
        slide_contents = get_comprehensive_yoga_content(main_asana, student_name)
        
        presentation = {
            "id": random.randint(1000, 9999),
            "title": f"{main_asana} - Comprehensive Yoga Guide",
            "asanas": [main_asana],
            "student_name": student_name,
            "generated_at": datetime.now().isoformat(),
            "slides": [],
            "status": "success",
            "message": "Created premium comprehensive presentation! 🎨"
        }
        
        # Create all 13 slides with enhanced content
        for i, content in enumerate(slide_contents):
            slide_data = {
                "title": f"Slide {i + 1}",
                "content": content,
                "image": None
            }
            
            if include_images:
                slide_data['image'] = get_canva_style_image(main_asana, i)
            
            presentation['slides'].append(slide_data)
        
        return jsonify(presentation)
    
    except Exception as e:
        return jsonify({
            'status': 'error',
            'message': f'Failed to generate: {str(e)}'
        }), 500

@app.route('/asanas', methods=['GET'])
def get_asanas():
    return jsonify({
        "asanas": YOGA_ASANAS,
        "total": len(YOGA_ASANAS),
        "message": "Premium comprehensive yoga presentations"
    })

@app.route('/download-pptx', methods=['POST'])
def download_pptx():
    try:
        data = request.get_json()
        presentation_data = data.get('presentation')
        
        if not presentation_data:
            return jsonify({'status': 'error', 'message': 'No presentation data'}), 400
        
        pptx_stream = create_canva_style_powerpoint(presentation_data)
        filename = f"comprehensive_yoga_{presentation_data['asanas'][0].replace(' ', '_').lower()}"
        if presentation_data.get('student_name'):
            filename += f"_{presentation_data['student_name'].replace(' ', '_').lower()}"
        filename += ".pptx"
        
        return send_file(
            pptx_stream,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        return jsonify({'status': 'error', 'message': f'Download failed: {str(e)}'}), 500

@app.route('/health', methods=['GET'])
def health_check():
    return jsonify({
        "status": "comprehensive! 🎨", 
        "message": "Enhanced Yoga Generator Running!",
        "features": [
            "5-6 Detailed Bullet Points Per Slide",
            "Enhanced Canva Styling & Typography", 
            "Professional Image Placement",
            "Comprehensive Content Coverage",
            "Geometric Design Elements"
        ]
    })

if __name__ == '__main__':
    print("🎨 STARTING COMPREHENSIVE YOGA GENERATOR...")
    print("=" * 65)
    print("✨ ENHANCED FEATURES:")
    print("• 5-6 Detailed Bullet Points Per Slide (Comprehensive Content)")
    print("• Enhanced Canva Styling with Professional Typography")
    print("• Geometric Design Elements & Gradient Backgrounds") 
    print("• Professional Image Placement with Enhanced Frames")
    print("• Complete 13-Slide Coverage with No Empty Spaces")
    print("=" * 65)
    app.run(debug=True, port=5000, host='127.0.0.1')